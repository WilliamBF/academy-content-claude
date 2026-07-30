# -*- coding: utf-8 -*-
"""Extract text, tables, notes, and embedded images from resource documents.

Supports: .pptx (python-pptx), .pdf (pdfplumber + pypdf fallback), .docx
(python-docx, pandoc fallback), .md / .txt (read directly).

Outputs into <out_dir>:
  extracted.json        - one entry per source file: ordered text/table/notes
                          blocks and a list of extracted image records
  assets/<file-slug>/   - extracted image binaries (png/jpg/...)
  images_manifest.json  - flat list of every image (path, source, bytes)

Usage:
  python extract_resources.py <out_dir> <file1> [file2 ...]

Images are saved locally. To use them in Thought Industries, upload them via
the TI editor or host at a URL. See /convert-course-to-html for the image
upload workflow.
"""
import sys, os, json, re


def slug(s):
    return re.sub(r"[^a-z0-9]+", "-",
                  os.path.splitext(os.path.basename(s))[0].lower()).strip("-") or "file"


# ---------- PPTX ----------

def _iter_pptx_shapes(shapes):
    """Yield every shape, descending into group shapes so nested images are not missed."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in _iter_pptx_shapes(shape.shapes):
                yield sub
        else:
            yield shape


def _pptx_image(shape):
    try:
        return shape.image
    except Exception:
        return None


def extract_pptx(path, asset_dir):
    from pptx import Presentation
    prs = Presentation(path)
    blocks, images = [], []
    for i, slide in enumerate(prs.slides, 1):
        blocks.append({"slide": i, "kind": "slide_start"})
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.has_text_frame:
                txt = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if txt.strip():
                    blocks.append({"slide": i, "kind": "text", "text": txt})
            if shape.has_table:
                rows = [" | ".join(c.text for c in r.cells) for r in shape.table.rows]
                blocks.append({"slide": i, "kind": "table", "text": "\n".join(rows)})
            img = _pptx_image(shape)
            if img is not None:
                fn = f"slide{i:02d}-{len(images)+1}.{img.ext}"
                fp = os.path.join(asset_dir, fn)
                with open(fp, "wb") as f:
                    f.write(img.blob)
                images.append({"slide": i, "file": fp, "bytes": len(img.blob),
                               "content_type": img.content_type})
        if slide.has_notes_slide:
            n = slide.notes_slide.notes_text_frame.text
            if n.strip():
                blocks.append({"slide": i, "kind": "notes", "text": n})
    return blocks, images


# ---------- PDF ----------

def extract_pdf(path, asset_dir):
    blocks, images = [], []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                t = page.extract_text() or ""
                if t.strip():
                    blocks.append({"page": i, "kind": "text", "text": t})
                for j, im in enumerate(page.images, 1):
                    try:
                        x0, top, x1, bottom = im["x0"], im["top"], im["x1"], im["bottom"]
                        crop = page.within_bbox((x0, top, x1, bottom)).to_image(resolution=150)
                        fn = f"page{i:02d}-{j}.png"
                        fp = os.path.join(asset_dir, fn)
                        crop.save(fp)
                        images.append({"page": i, "file": fp})
                    except Exception:
                        pass
    except Exception:
        try:
            from pypdf import PdfReader
            r = PdfReader(path)
            for i, page in enumerate(r.pages, 1):
                t = page.extract_text() or ""
                if t.strip():
                    blocks.append({"page": i, "kind": "text", "text": t})
        except Exception as e:
            print(f"  WARNING: could not extract PDF text ({e})")
    return blocks, images


# ---------- DOCX ----------

def extract_docx(path, asset_dir):
    blocks, images = [], []
    try:
        import docx
        d = docx.Document(path)
        for p in d.paragraphs:
            if p.text.strip():
                blocks.append({"kind": "text", "text": p.text})
        for t in d.tables:
            rows = [" | ".join(c.text for c in r.cells) for r in t.rows]
            blocks.append({"kind": "table", "text": "\n".join(rows)})
        for rel in d.part.rels.values():
            if "image" in rel.reltype:
                try:
                    blob = rel.target_part.blob
                    ext = rel.target_part.partname.ext.lstrip(".") or "png"
                    fn = f"img-{len(images)+1}.{ext}"
                    fp = os.path.join(asset_dir, fn)
                    with open(fp, "wb") as f:
                        f.write(blob)
                    images.append({"file": fp, "bytes": len(blob)})
                except Exception:
                    pass
    except Exception:
        try:
            import subprocess
            md = subprocess.run(["pandoc", "-t", "markdown", path],
                                capture_output=True, text=True).stdout
            blocks.append({"kind": "text", "text": md})
        except Exception as e:
            print(f"  WARNING: could not extract DOCX ({e}); pandoc not available?")
    return blocks, images


# ---------- MD / TXT ----------

def extract_text_file(path, asset_dir):
    txt = open(path, encoding="utf-8", errors="ignore").read()
    return [{"kind": "text", "text": txt}], []


EXTRACTORS = {
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".md": extract_text_file,
    ".txt": extract_text_file,
}


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    out_dir = sys.argv[1]
    files = sys.argv[2:]
    os.makedirs(out_dir, exist_ok=True)
    assets_root = os.path.join(out_dir, "assets")
    os.makedirs(assets_root, exist_ok=True)

    resources, all_images = [], []
    for path in files:
        ext = os.path.splitext(path)[1].lower()
        fn = EXTRACTORS.get(ext)
        if not fn:
            print(f"  SKIP (unsupported format): {path}")
            continue
        adir = os.path.join(assets_root, slug(path))
        os.makedirs(adir, exist_ok=True)
        try:
            blocks, images = fn(path, adir)
        except Exception as e:
            print(f"  ERROR extracting {path}: {e}")
            blocks, images = [], []
        for im in images:
            im["source"] = os.path.basename(path)
        all_images.extend(images)
        resources.append({
            "name": os.path.basename(path),
            "ext": ext,
            "blocks": blocks,
            "images": images,
        })
        print(f"  {os.path.basename(path)}: {len(blocks)} blocks, {len(images)} images")

    json.dump({"resources": resources},
              open(os.path.join(out_dir, "extracted.json"), "w"),
              indent=2, ensure_ascii=False)
    json.dump(all_images,
              open(os.path.join(out_dir, "images_manifest.json"), "w"),
              indent=2, ensure_ascii=False)
    print(
        f"\nDone. Wrote extracted.json | "
        f"{len(resources)} file(s), {len(all_images)} image(s) total"
    )


if __name__ == "__main__":
    main()
